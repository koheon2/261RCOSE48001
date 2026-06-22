"""
ResearcherHub Codex — A0 Poster Generator
A0 = 841mm × 1189mm  →  pptx uses EMU: 1mm = 36000 EMU
"""

from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Color Palette (Dark mode matching the app) ──────────────────────────────
BG_DARK      = RGBColor(0x08, 0x0d, 0x14)   # #080d14
BG_CARD      = RGBColor(0x0d, 0x18, 0x27)   # #0d1827
BG_CARD2     = RGBColor(0x0f, 0x20, 0x35)   # #0f2035
TEAL         = RGBColor(0x0f, 0x76, 0x6e)   # #0f766e
TEAL_LIGHT   = RGBColor(0x14, 0xb8, 0xa6)   # #14b8a6
BLUE         = RGBColor(0x25, 0x63, 0xeb)   # #2563eb
BLUE_LIGHT   = RGBColor(0x60, 0xa5, 0xfa)   # #60a5fa
PURPLE       = RGBColor(0x7c, 0x3a, 0xed)   # #7c3aed
PURPLE_LIGHT = RGBColor(0xa7, 0x8b, 0xfa)   # #a78bfa
ORANGE       = RGBColor(0xf9, 0x73, 0x16)   # #f97316
YELLOW       = RGBColor(0xfb, 0xbf, 0x24)   # #fbbf24
WHITE        = RGBColor(0xff, 0xff, 0xff)
GRAY_300     = RGBColor(0xd1, 0xd5, 0xdb)
GRAY_400     = RGBColor(0x9c, 0xa3, 0xaf)
GRAY_500     = RGBColor(0x6b, 0x72, 0x80)
GRAY_700     = RGBColor(0x37, 0x41, 0x51)

# ── Dimensions ───────────────────────────────────────────────────────────────
W_MM = 841
H_MM = 1189
MARGIN = 22  # mm

prs = Presentation()
prs.slide_width  = Mm(W_MM)
prs.slide_height = Mm(H_MM)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)


# ── Helpers ──────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Mm
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE
        Mm(x), Mm(y), Mm(w), Mm(h)
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=Pt(12), bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, italic=False,
                line_spacing=None, word_wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as Pt2
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline_textbox(slide, x, y, w, h, lines, font_size=Pt(11),
                          color=GRAY_300, align=PP_ALIGN.LEFT,
                          bold=False, font_name="Calibri",
                          line_spacing_pt=14):
    txBox = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        # line spacing
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing_pt * 100)))
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox


def add_labeled_box(slide, x, y, w, h, label, value, label_color=TEAL_LIGHT,
                    value_color=WHITE, bg=BG_CARD2, border=TEAL,
                    label_size=Pt(9), value_size=Pt(22)):
    add_rect(slide, x, y, w, h, fill_color=bg, line_color=border, line_width=Pt(1.2))
    # label
    add_textbox(slide, x+2, y+2.5, w-4, 10, label,
                font_size=label_size, color=label_color, align=PP_ALIGN.CENTER, bold=True)
    # value
    add_textbox(slide, x+2, y+10, w-4, h-14, value,
                font_size=value_size, color=value_color, align=PP_ALIGN.CENTER, bold=True)


def add_section_header(slide, x, y, w, label, accent=TEAL):
    # accent line
    add_rect(slide, x, y, 4, 8, fill_color=accent)
    add_textbox(slide, x+6, y-1, w-10, 11, label,
                font_size=Pt(13), bold=True, color=WHITE)
    add_rect(slide, x, y+9, w, 0.5, fill_color=accent)


# ═══════════════════════════════════════════════════════════════════════════
# 1. BACKGROUND
# ═══════════════════════════════════════════════════════════════════════════
add_rect(slide, 0, 0, W_MM, H_MM, fill_color=BG_DARK)

# Subtle grid lines decoration
for xi in range(0, W_MM, 40):
    add_rect(slide, xi, 0, 0.3, H_MM,
             fill_color=RGBColor(0x14, 0x25, 0x38))
for yi in range(0, H_MM, 40):
    add_rect(slide, 0, yi, W_MM, 0.3,
             fill_color=RGBColor(0x14, 0x25, 0x38))

# ═══════════════════════════════════════════════════════════════════════════
# 2. HEADER BAND
# ═══════════════════════════════════════════════════════════════════════════
add_rect(slide, 0, 0, W_MM, 62, fill_color=RGBColor(0x05, 0x10, 0x1e))
# teal accent bar
add_rect(slide, 0, 62, W_MM, 2.5, fill_color=TEAL)
# gradient-like right decoration
add_rect(slide, W_MM-160, 0, 160, 64,
         fill_color=RGBColor(0x0a, 0x1c, 0x30))

# Logo / Icon block
add_rect(slide, MARGIN, 10, 20, 20, fill_color=TEAL,
         line_color=TEAL_LIGHT, line_width=Pt(0.5))
add_textbox(slide, MARGIN+1, 11, 18, 18, "RH",
            font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Title
add_textbox(slide, MARGIN+25, 7, 550, 22, "ResearcherHub Codex",
            font_size=Pt(36), bold=True, color=WHITE, font_name="Calibri")
add_textbox(slide, MARGIN+25, 31, 560, 14,
            "AI/CS Research Intelligence Platform  ·  Publication-Time Attribution  ·  23M Papers  ·  3.2M Researchers",
            font_size=Pt(13), color=TEAL_LIGHT, font_name="Calibri")

# Right info
add_textbox(slide, W_MM-155, 10, 140, 14,
            "Koheon2  ·  2026",
            font_size=Pt(11), color=GRAY_400, align=PP_ALIGN.RIGHT)
add_textbox(slide, W_MM-155, 24, 140, 14,
            "ResearcherHub Codex",
            font_size=Pt(10), color=GRAY_500, align=PP_ALIGN.RIGHT)
add_textbox(slide, W_MM-155, 38, 140, 16,
            "codex/topic-lineage-mvp",
            font_size=Pt(9), color=TEAL, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# 3. KEY STATS ROW
# ═══════════════════════════════════════════════════════════════════════════
Y_STATS = 70
stats = [
    ("Total Papers",     "23M+"),
    ("Researchers",      "3.2M"),
    ("Pub-Time Affils.", "21.4M"),
    ("Citation Edges",   "1.87M"),
    ("Countries",        "190+"),
    ("Institutions",     "63,985"),
    ("Intent Types",     "10"),
    ("Quality Flags",    "16.7M"),
]
cols = len(stats)
box_w = (W_MM - 2*MARGIN - (cols-1)*3) / cols
for i, (lbl, val) in enumerate(stats):
    bx = MARGIN + i * (box_w + 3)
    add_labeled_box(slide, bx, Y_STATS, box_w, 28,
                    lbl, val,
                    label_color=GRAY_400, value_color=TEAL_LIGHT,
                    bg=BG_CARD, border=GRAY_700,
                    label_size=Pt(8), value_size=Pt(18))


# ═══════════════════════════════════════════════════════════════════════════
# 4. MAIN CONTENT — 3 column layout
# ═══════════════════════════════════════════════════════════════════════════
Y_MAIN = 106   # below stats row
COL_W  = 255
COL_GAP = 12
COL1_X = MARGIN
COL2_X = MARGIN + COL_W + COL_GAP
COL3_X = MARGIN + 2*(COL_W + COL_GAP)

# ── COLUMN 1 ── OVERVIEW & MOTIVATION ──────────────────────────────────────
y = Y_MAIN
add_section_header(slide, COL1_X, y, COL_W, "Overview & Motivation", TEAL)
y += 14

add_rect(slide, COL1_X, y, COL_W, 110, fill_color=BG_CARD,
         line_color=GRAY_700, line_width=Pt(0.5))

overview_text = (
    "ResearcherHub Codex is a full-stack research visualization platform "
    "designed to reveal the true landscape of AI and computer science scholarship "
    "across 23 million papers from OpenAlex.\n\n"
    "Unlike conventional research databases that rely on a researcher's current "
    "affiliation, ResearcherHub tracks publication-time affiliations — capturing "
    "where an author actually was when each paper was written. This prevents "
    "researcher mobility from distorting country and institution contribution metrics.\n\n"
    "The platform combines pre-computed summary analytics, interactive 3D globe "
    "visualization, natural language query routing, and citation lineage graphs to "
    "provide a comprehensive view of who contributed what, from where, and when."
)
add_textbox(slide, COL1_X+4, y+4, COL_W-8, 102, overview_text,
            font_size=Pt(10.5), color=GRAY_300, line_spacing=13.5)

y += 115

# Problem Statement box
add_rect(slide, COL1_X, y, COL_W, 40, fill_color=RGBColor(0x0a,0x1f,0x12),
         line_color=TEAL, line_width=Pt(1.5))
add_textbox(slide, COL1_X+4, y+3, COL_W-8, 8, "KEY PROBLEM",
            font_size=Pt(8), bold=True, color=TEAL_LIGHT)
add_textbox(slide, COL1_X+4, y+12, COL_W-8, 26,
            "Existing databases attribute papers to a researcher's current institution. "
            "ResearcherHub anchors each paper to the affiliation at publication time — "
            "enabling accurate retrospective analysis of national and institutional research output.",
            font_size=Pt(10), color=GRAY_300, line_spacing=13)
y += 45

add_section_header(slide, COL1_X, y, COL_W, "Architecture", BLUE)
y += 14

# Architecture flow boxes
arch_layers = [
    (BLUE,         "Data Layer  (PostgreSQL 16)",
     "23M papers · 21.4M pub-time affiliations\nPre-computed summary tables (no live joins)\nCitation edges · Quality flags · Facets"),
    (PURPLE,       "API Layer  (FastAPI + asyncpg)",
     "9 route modules · OpenAI intent parsing\nSummary-table-first queries · Graph endpoints\nAPScheduler 6-hour pipeline runs"),
    (TEAL,         "Frontend Layer  (React 19 + Vite 8)",
     "13 page components · Three.js / Cesium globe\nD3 citation graphs · Real-time polling\nDark-mode visual design system"),
]
for accent, title, body in arch_layers:
    add_rect(slide, COL1_X, y, COL_W, 36, fill_color=BG_CARD2,
             line_color=accent, line_width=Pt(1.5))
    add_textbox(slide, COL1_X+5, y+4, COL_W-10, 9, title,
                font_size=Pt(10), bold=True, color=accent)
    add_textbox(slide, COL1_X+5, y+14, COL_W-10, 22, body,
                font_size=Pt(9.5), color=GRAY_300, line_spacing=12.5)
    # arrow
    if accent != TEAL:
        add_textbox(slide, COL1_X + COL_W//2 - 5, y+36, 10, 5, "▼",
                    font_size=Pt(9), color=GRAY_500, align=PP_ALIGN.CENTER)
    y += 40

y += 5
add_section_header(slide, COL1_X, y, COL_W, "Tech Stack", PURPLE)
y += 14

tech = [
    ("Frontend",  "React 19 · TypeScript 5.9 · Vite 8", BLUE_LIGHT),
    ("3D / Maps", "Three.js · Cesium · React Three Fiber", TEAL_LIGHT),
    ("Backend",   "FastAPI · SQLAlchemy 2.0 · asyncpg", PURPLE_LIGHT),
    ("Database",  "PostgreSQL 16 · Alembic migrations", ORANGE),
    ("LLM",       "OpenAI GPT-4o-mini (intent routing)", YELLOW),
    ("Pipeline",  "APScheduler · boto3 · httpx", GRAY_300),
]
for label, val, col in tech:
    add_rect(slide, COL1_X, y, COL_W, 12, fill_color=BG_CARD,
             line_color=GRAY_700, line_width=Pt(0.3))
    add_textbox(slide, COL1_X+4, y+1.5, 52, 9, label,
                font_size=Pt(8.5), bold=True, color=col)
    add_textbox(slide, COL1_X+58, y+1.5, COL_W-62, 9, val,
                font_size=Pt(8.5), color=GRAY_300)
    y += 13


# ── COLUMN 2 ── FEATURES & PAGES ───────────────────────────────────────────
y = Y_MAIN
add_section_header(slide, COL2_X, y, COL_W, "Platform Features & Pages", TEAL)
y += 14

pages = [
    (TEAL,         "Globe Explorer",
     "3D/2D Cesium globe showing researcher point clouds\nglobally. Filter by field, country, city."),
    (BLUE,         "Topic Lineage Graph",
     "Multi-seed citation lineage graph. Reveals intellectual\nancestors of research topics using multi_seed_v0 scoring."),
    (PURPLE,       "Trending Topics",
     "Facet-driven trending engine ranks topics by growth %\nacross 4 axes: Aboutness · Method · Task · Application."),
    (ORANGE,       "Leaderboard",
     "Rankings by country, institution, researcher. Sortable\nby citations, contributions, hotness, papers, FWCI."),
    (YELLOW,       "Progress Tracker",
     "Year-over-year trend charts showing contribution and\ncitation growth for countries and fields."),
    (TEAL_LIGHT,   "Paper Detail",
     "Full paper view: metadata, authors, affiliations,\nquality flags, facets, reference summaries."),
    (BLUE_LIGHT,   "Paper Graph",
     "Single-paper citation lineage: seed → references →\n2-hop prerequisites → related works."),
    (PURPLE_LIGHT, "Institution Profile",
     "Deep-dive: top fields, key researchers, trend data,\nand representative papers for any institution."),
    (GRAY_300,     "Paper Timeline",
     "Topic-based paper discovery sorted by impact,\norganized by temporal research trends."),
    (TEAL,         "Researcher DNA",
     "Individual researcher profile: contribution history,\nco-author network, topic evolution."),
    (BLUE,         "Research Map",
     "Concept map with timeline evolution, research camps,\nkey papers, and narrative analysis."),
    (ORANGE,       "Benchmark Page",
     "Benchmark / dataset tracking and metric extraction\nfrom papers (MVP stage)."),
]

for accent, name, desc in pages:
    add_rect(slide, COL2_X, y, COL_W, 28, fill_color=BG_CARD,
             line_color=accent, line_width=Pt(1.0))
    # colored left bar
    add_rect(slide, COL2_X, y, 3.5, 28, fill_color=accent)
    add_textbox(slide, COL2_X+7, y+3, COL_W-11, 8, name,
                font_size=Pt(10), bold=True, color=accent)
    add_textbox(slide, COL2_X+7, y+12, COL_W-11, 16, desc,
                font_size=Pt(9), color=GRAY_300, line_spacing=12)
    y += 30

# Add universal search feature
y += 2
add_rect(slide, COL2_X, y, COL_W, 32, fill_color=RGBColor(0x10,0x1e,0x38),
         line_color=BLUE, line_width=Pt(1.5))
add_textbox(slide, COL2_X+5, y+4, COL_W-10, 9,
            "Universal Natural Language Search",
            font_size=Pt(10.5), bold=True, color=BLUE_LIGHT)
add_textbox(slide, COL2_X+5, y+14, COL_W-10, 18,
            "OpenAI GPT-4o-mini parses free-text queries into 10 structured intent types "
            "(comparison, progress, leaderboard, trending, institution_profile, researcher_search, "
            "paper_detail, paper_graph, topic_lineage, paper_timeline) and routes to the "
            "corresponding visualization page. Bilingual: Korean + English.",
            font_size=Pt(9), color=GRAY_300, line_spacing=12)


# ── COLUMN 3 ── NOVEL CONTRIBUTIONS & METHODS ─────────────────────────────
y = Y_MAIN
add_section_header(slide, COL3_X, y, COL_W, "Novel Contributions", ORANGE)
y += 14

contributions = [
    (TEAL, "Publication-Time Affiliation Attribution",
     [
         "21.4M paper-author-affiliation tuples anchored to publish time.",
         "Prevents researcher mobility from distorting country/institution stats.",
         "Confidence scoring + ROR canonical matching (99.6% coverage).",
         "Enables: 'What was MIT's true AI contribution in 2020?'",
     ]),
    (BLUE, "Weak-Label Facet Taxonomy",
     [
         "4-axis paper classification: Aboutness · Method · Task · Application.",
         "Not just topic tags — enables deep slicing (e.g., 'Transformer in medical imaging').",
         "Rule/keyword-based v0; embedding/LLM taxonomy planned.",
         "Enables facet-driven trending, timeline discovery, and lineage scoring.",
     ]),
    (PURPLE, "Multi-Seed Citation Lineage Scoring  (multi_seed_v0)",
     [
         "Score = α·seed_reach + β·edge_count + γ·citations + δ·depth_boost − type_penalty",
         "seed_reach: How many seed papers reference/cite this ancestor?",
         "edge_count: How often does it appear across reference chains?",
         "Reveals intellectual foundations of a field, not just direct citations.",
     ]),
    (ORANGE, "Quality-Filtered Product Layer",
     [
         "16.7M quality flags across 23M papers; conservative_v0 policy.",
         "Original data preserved — no destructive edits.",
         "Summary tables respect filter; UI can toggle policy.",
         "Enables A/B testing of different quality thresholds.",
     ]),
    (YELLOW, "Summary Table Architecture",
     [
         "Pre-computed aggregates: country_stats, institution_field_stats, facet_year_summary, etc.",
         "All API endpoints query summary tables — no live joins on 23M rows.",
         "Nightly refresh via APScheduler; resumable/idempotent pipeline.",
         "Enables sub-100ms API responses across the full dataset.",
     ]),
]

for accent, title, bullets in contributions:
    nh = 14 + len(bullets)*12 + 4
    add_rect(slide, COL3_X, y, COL_W, nh, fill_color=BG_CARD2,
             line_color=accent, line_width=Pt(1.5))
    add_textbox(slide, COL3_X+5, y+4, COL_W-10, 9, title,
                font_size=Pt(10.5), bold=True, color=accent)
    by = y + 14
    for bullet in bullets:
        add_textbox(slide, COL3_X+8, by, COL_W-12, 11, "• " + bullet,
                    font_size=Pt(9.5), color=GRAY_300, line_spacing=11.5)
        by += 11.5
    y += nh + 5


# ═══════════════════════════════════════════════════════════════════════════
# 5. DATA PIPELINE SECTION (full width, bottom)
# ═══════════════════════════════════════════════════════════════════════════
Y_PIPE = 1090
add_rect(slide, 0, Y_PIPE-2, W_MM, 2, fill_color=TEAL)
add_rect(slide, 0, Y_PIPE, W_MM, H_MM - Y_PIPE, fill_color=RGBColor(0x05,0x10,0x1e))

add_section_header(slide, MARGIN, Y_PIPE+4, W_MM-2*MARGIN,
                   "Data Pipeline & Processing", TEAL)

PIPE_Y = Y_PIPE + 18
pipe_steps = [
    (TEAL,        "OpenAlex Snapshot\nIngest",
     "23M papers\n3.2M researchers\nRaw JSON → PostgreSQL"),
    (BLUE,        "Publication-Time\nAffiliation Backfill",
     "21.4M rows\nROR normalization\nFuzzy + exact match"),
    (PURPLE,      "Facet\nClassification",
     "Keyword/regex\n4-axis taxonomy\npaper_facets table"),
    (ORANGE,      "Quality\nFiltering",
     "conservative_v0\n16.7M flags\nSummary-safe"),
    (YELLOW,      "Citation\nEnrichment",
     "Top 30K papers\nOpenAlex API\n1.87M edges"),
    (TEAL_LIGHT,  "Summary Table\nRefresh",
     "Nightly batch\nPre-aggregation\n<100ms API"),
    (BLUE_LIGHT,  "Lineage Graph\nComputation",
     "multi_seed_v0\nOn-demand (MVP)\nCache layer planned"),
]

pipe_w = (W_MM - 2*MARGIN - 6*8) / 7
for i, (accent, step_title, step_body) in enumerate(pipe_steps):
    px = MARGIN + i*(pipe_w + 8)
    add_rect(slide, px, PIPE_Y, pipe_w, 72, fill_color=BG_CARD,
             line_color=accent, line_width=Pt(1.5))
    # step number
    add_rect(slide, px+pipe_w/2-5, PIPE_Y+2, 10, 10, fill_color=accent)
    add_textbox(slide, px+pipe_w/2-5, PIPE_Y+2.5, 10, 8, str(i+1),
                font_size=Pt(7), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, px+2, PIPE_Y+14, pipe_w-4, 16, step_title,
                font_size=Pt(9), bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_textbox(slide, px+2, PIPE_Y+32, pipe_w-4, 36, step_body,
                font_size=Pt(8.5), color=GRAY_300, align=PP_ALIGN.CENTER,
                line_spacing=11)
    if i < len(pipe_steps)-1:
        add_textbox(slide, px+pipe_w+1, PIPE_Y+32, 8, 8, "→",
                    font_size=Pt(10), bold=True, color=GRAY_500,
                    align=PP_ALIGN.CENTER)


# Footer
add_rect(slide, 0, H_MM-10, W_MM, 10, fill_color=RGBColor(0x03,0x09,0x0f))
add_textbox(slide, MARGIN, H_MM-9, 300, 8,
            "ResearcherHub Codex  ·  Open-source research intelligence platform",
            font_size=Pt(7), color=GRAY_500)
add_textbox(slide, W_MM//2-100, H_MM-9, 200, 8,
            "OpenAlex · PostgreSQL 16 · FastAPI · React 19 · Three.js · Cesium",
            font_size=Pt(7), color=GRAY_500, align=PP_ALIGN.CENTER)
add_textbox(slide, W_MM-MARGIN-180, H_MM-9, 180, 8,
            "github: koheon2/researcherhub-codex  ·  2026",
            font_size=Pt(7), color=GRAY_500, align=PP_ALIGN.RIGHT)


# ── SAVE ────────────────────────────────────────────────────────────────────
out_path = "/Users/amadeus/Desktop/researcherhub-codex/ResearcherHub_Codex_A0_Poster.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
